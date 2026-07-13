from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from .config import (
    OCR_DPI,
    OCR_ENABLED,
    OCR_LANG,
    OCR_MIN_DOCUMENT_CHARS,
    OCR_MAX_PAGES,
    OCR_MIN_PAGE_CHARS,
    OCR_TESSDATA_DIR,
    PDF_TEXT_ENGINE,
    TESSERACT_CMD,
)


SUPPORTED_SUFFIXES = {".pdf", ".pptx", ".docx", ".txt", ".md"}


def extract_pages(path: Path) -> list[tuple[int | None, str]]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix in {".txt", ".md"}:
        return [(None, path.read_text(encoding="utf-8", errors="ignore"))]
    raise ValueError(f"暂不支持的文件类型：{suffix}")


def extract_pdf(path: Path) -> list[tuple[int | None, str]]:
    pages = extract_pdf_text(path)
    if should_ocr_document(pages):
        ocr_candidates = [
            int(page_number)
            for page_number, text in pages
            if page_number is not None and should_ocr_page(text)
        ][:OCR_MAX_PAGES]
        if ocr_candidates:
            ocr_pages = ocr_pdf_pages(path, ocr_candidates)
            pages = [
                (page_number, ocr_pages.get(page_number, text) or text)
                for page_number, text in pages
            ]
    return pages


def extract_pdf_text(path: Path) -> list[tuple[int | None, str]]:
    if PDF_TEXT_ENGINE in {"pymupdf", "auto"}:
        try:
            return extract_pdf_text_pymupdf(path)
        except Exception:
            if PDF_TEXT_ENGINE == "pymupdf":
                raise
    return extract_pdf_text_pypdf(path)


def extract_pdf_text_pymupdf(path: Path) -> list[tuple[int | None, str]]:
    import fitz

    pages: list[tuple[int | None, str]] = []
    with fitz.open(str(path)) as document:
        for index, page in enumerate(document, start=1):
            pages.append((index, clean_text(page.get_text("text") or "")))
    return pages


def extract_pdf_text_pypdf(path: Path) -> list[tuple[int | None, str]]:
    reader = PdfReader(str(path))
    pages: list[tuple[int | None, str]] = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append((index, clean_text(text)))
    return pages


def should_ocr_document(pages: list[tuple[int | None, str]]) -> bool:
    if not OCR_ENABLED:
        return False
    total_chars = sum(len(text.strip()) for _, text in pages)
    return total_chars < OCR_MIN_DOCUMENT_CHARS


def should_ocr_page(text: str) -> bool:
    return len(text.strip()) < OCR_MIN_PAGE_CHARS


def ocr_pdf_pages(path: Path, page_numbers: list[int]) -> dict[int, str]:
    try:
        import fitz
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        missing = exc.name or "OCR dependency"
        raise RuntimeError(
            f"Local OCR needs optional dependency '{missing}'. Run `pip install -r requirements.txt`."
        ) from exc

    if TESSERACT_CMD:
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

    zoom = OCR_DPI / 72
    matrix = fitz.Matrix(zoom, zoom)
    results: dict[int, str] = {}
    try:
        with fitz.open(str(path)) as document:
            for page_number in page_numbers:
                page = document.load_page(page_number - 1)
                pixmap = page.get_pixmap(matrix=matrix, alpha=False)
                image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
                config = f'--tessdata-dir "{OCR_TESSDATA_DIR}"' if OCR_TESSDATA_DIR else ""
                text = pytesseract.image_to_string(image, lang=OCR_LANG, config=config)
                results[page_number] = clean_text(text)
    except pytesseract.TesseractNotFoundError as exc:
        raise RuntimeError(
            "Tesseract OCR is not installed or not on PATH. Install Tesseract, or set TESSERACT_CMD in .env."
        ) from exc
    except pytesseract.TesseractError as exc:
        raise RuntimeError(f"Tesseract OCR failed for {path.name}: {exc}") from exc
    return results


def extract_pptx(path: Path) -> list[tuple[int | None, str]]:
    presentation = Presentation(str(path))
    pages: list[tuple[int | None, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
        pages.append((index, clean_text("\n".join(parts))))
    return pages


def extract_docx(path: Path) -> list[tuple[int | None, str]]:
    doc = Document(str(path))
    parts = [paragraph.text for paragraph in doc.paragraphs]
    return [(None, clean_text("\n".join(parts)))]


def clean_text(text: str) -> str:
    lines = [line.strip() for line in text.replace("\x00", " ").splitlines()]
    return "\n".join(line for line in lines if line)


def chunk_pages(
    pages: list[tuple[int | None, str]],
    max_chars: int = 1800,
    overlap: int = 200,
) -> list[tuple[int | None, int, str]]:
    chunks: list[tuple[int | None, int, str]] = []
    chunk_index = 0
    for page, text in pages:
        if not text.strip():
            continue
        start = 0
        while start < len(text):
            end = min(start + max_chars, len(text))
            chunk = text[start:end].strip()
            if chunk:
                chunks.append((page, chunk_index, chunk))
                chunk_index += 1
            if end >= len(text):
                break
            start = max(0, end - overlap)
    return chunks


def sample_for_llm(pages: list[tuple[int | None, str]], max_chars: int = 18000) -> str:
    selected: list[str] = []
    total = 0
    for page, text in pages:
        if not text.strip():
            continue
        prefix = f"\n\n[第{page}页]\n" if page else "\n\n[文档]\n"
        block = prefix + text.strip()
        if total + len(block) > max_chars:
            selected.append(block[: max_chars - total])
            break
        selected.append(block)
        total += len(block)
    return "".join(selected).strip()

